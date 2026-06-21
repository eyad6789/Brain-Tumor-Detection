# -*- coding: utf-8 -*-
"""Render NeuroScan_EN.pptx and NeuroScan_AR.pptx from slides_data.SLIDES.

Everything is drawn on blank layouts for full control of the dark
"reading-room" theme. Arabic slides set paragraph rtl + right alignment and a
complex-script (cs) font so Arabic shapes correctly while embedded Latin tech
terms (ResNet50, FastAPI…) stay in a clean Latin face.

Run:  ../.venv/bin/python build_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

import theme as T
from i18n import tr
from slides_data import SLIDES

HERE = os.path.dirname(__file__)
ASSETS = os.path.join(HERE, "..", "assets")
OUT = os.path.join(HERE, "..")

EMU_IN = 914400
SW, SH = 13.333, 7.5            # 16:9 inches
ML = 0.85                       # left/right margin

def C(hexstr): return RGBColor.from_string(hexstr)

CYAN, VIO, GRN, AMB = C(T.PRIMARY), C(T.ACCENT), C(T.SAFE), C(T.ALERT)
INK, INKS, MUT = C(T.INK), C(T.INK_STRONG), C(T.MUTED)
BG, SURF, SURF2, LINE = C(T.BG), C(T.SURFACE), C(T.SURFACE_2), C(T.LINE)


def _set_cs(run, typeface):
    """Set the complex-script (Arabic) typeface on a run."""
    rPr = run._r.get_or_add_rPr()
    for el in rPr.findall(qn("a:cs")):
        rPr.remove(el)
    cs = rPr.makeelement(qn("a:cs"), {"typeface": typeface})
    rPr.append(cs)


def _set_rtl(paragraph, rtl):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1" if rtl else "0")


class Deck:
    def __init__(self, lang):
        self.lang = lang
        self.rtl = lang == "ar"
        self.latin = T.FONT_DISPLAY                  # Calibri for Latin
        self.arface = T.FONT_AR if self.rtl else T.FONT_DISPLAY
        self.prs = Presentation()
        self.prs.slide_width = Emu(int(SW * EMU_IN))
        self.prs.slide_height = Emu(int(SH * EMU_IN))
        self.blank = self.prs.slide_layouts[6]

    # ── primitives ───────────────────────────────────────────────────────
    def slide(self):
        s = self.prs.slides.add_slide(self.blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BG
        return s

    def rect(self, s, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
        sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
        sp.shadow.inherit = False
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line; sp.line.width = Pt(lw)
        return sp

    def tb(self, s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
        box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = Inches(0.02)
        return tf

    def para(self, tf, runs, *, size, color=INK, bold=False, align=None,
             first=False, before=6, after=2, line=1.06, mono=False):
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
        if align is None:
            align = PP_ALIGN.RIGHT if self.rtl else PP_ALIGN.LEFT
        p.alignment = align
        _set_rtl(p, self.rtl)
        p.space_before = Pt(before); p.space_after = Pt(after)
        try: p.line_spacing = line
        except Exception: pass
        if isinstance(runs, str):
            runs = [(runs, color, bold, size)]
        for spec in runs:
            text, rc, rb, rs = spec
            r = p.add_run(); r.text = text
            r.font.size = Pt(rs); r.font.bold = rb; r.font.color.rgb = rc
            r.font.name = T.FONT_MONO if mono else self.latin
            if self.rtl and not mono:
                _set_cs(r, self.arface)
        return p

    def accent_bar(self, s, l, t, w=1.4, color=CYAN):
        self.rect(s, l, t, w, 0.07, fill=color)

    def heading(self, s, text, *, top=0.62):
        tf = self.tb(s, ML, top, SW - 2 * ML, 1.0)
        self.para(tf, [(text, INKS, True, 30)], size=30, first=True, before=0, after=4)
        self.accent_bar(s, SW - ML - 1.4 if self.rtl else ML, top + 0.74)

    def caption(self, s, text, *, top):
        tf = self.tb(s, ML, top, SW - 2 * ML, 0.9)
        self.para(tf, [(text, MUT, False, 13.5)], size=13.5, first=True,
                  before=0, after=0, line=1.15)

    def fit_image(self, s, path, box):
        l, t, w, h = box
        iw, ih = Image.open(path).size
        ar = iw / ih
        if w / h > ar:
            nw, nh = h * ar, h
        else:
            nw, nh = w, w / ar
        nl, nt = l + (w - nw) / 2, t + (h - nh) / 2
        s.shapes.add_picture(path, Inches(nl), Inches(nt), Inches(nw), Inches(nh))

    def bullet_block(self, s, items, *, top, width=None, size=18, gap=10, color=INK):
        width = width or (SW - 2 * ML)
        left = SW - ML - width if self.rtl else ML
        tf = self.tb(s, left, top, width, SH - top - 0.5)
        for i, it in enumerate(items):
            txt = tr(it, self.lang)
            runs = ([("•  ", CYAN, True, size), (txt, color, False, size)] if not self.rtl
                    else [(txt, color, False, size), ("  •", CYAN, True, size)])
            # For RTL the bullet run must be first logically to sit on the right:
            if self.rtl:
                runs = [("•  ", CYAN, True, size), (txt, color, False, size)]
            self.para(tf, runs, size=size, before=(0 if i == 0 else gap),
                      after=0, line=1.12, first=(i == 0))
        return tf

    # ── slide kinds ───────────────────────────────────────────────────────
    def s_title(self, d):
        s = self.slide()
        # backdrop accents
        self.rect(s, -1, 5.7, 15, 3, fill=SURF)
        self.accent_bar(s, (SW - ML - 2.0) if self.rtl else ML, 2.05, w=2.0, color=CYAN)
        tf = self.tb(s, ML, 1.4, SW - 2 * ML, 0.6)
        self.para(tf, [(tr(d["tag"], self.lang).upper() if self.lang == "en" else tr(d["tag"], self.lang),
                        CYAN, True, 14)], size=14, first=True, before=0, after=0)
        tf = self.tb(s, ML, 2.25, SW - 2 * ML, 1.5)
        self.para(tf, [(tr(d["title"], self.lang), INKS, True, 60)], size=60, first=True, before=0, after=0)
        tf = self.tb(s, ML, 3.7, SW - 2 * ML, 1.6)
        self.para(tf, [(tr(d["subtitle"], self.lang), INK, False, 22)], size=22, first=True,
                  before=0, after=0, line=1.2)
        tf = self.tb(s, ML, 6.55, SW - 2 * ML, 0.6)
        self.para(tf, [(tr(d["foot"], self.lang), MUT, False, 14)], size=14, first=True, before=0, after=0)

    def s_section(self, d):
        s = self.slide()
        tf = self.tb(s, ML, 1.0, SW - 2 * ML, 2.0)
        self.para(tf, [(d["num"], LINE, True, 150)], size=150, first=True, before=0, after=0,
                  align=(PP_ALIGN.RIGHT if self.rtl else PP_ALIGN.LEFT))
        self.accent_bar(s, SW - ML - 2.0 if self.rtl else ML, 3.85, w=2.0, color=CYAN)
        tf = self.tb(s, ML, 4.0, SW - 2 * ML, 1.4)
        self.para(tf, [(tr(d["title"], self.lang), INKS, True, 40)], size=40, first=True, before=0, after=4)
        self.para(tf, [(tr(d["subtitle"], self.lang), MUT, False, 19)], size=19, before=4, after=0)

    def s_bullets(self, d):
        s = self.slide()
        self.heading(s, tr(d["title"], self.lang))
        top = 1.75
        if d.get("lead"):
            tf = self.tb(s, ML, top, SW - 2 * ML, 0.8)
            self.para(tf, [(tr(d["lead"], self.lang), INK, False, 19)], size=19, first=True,
                      before=0, after=0, line=1.15)
            top += 0.72
        self.bullet_block(s, d["bullets"], top=top)
        if d.get("foot"):
            self.rect(s, ML, 6.62, SW - 2 * ML, 0.02, fill=LINE)
            self.caption(s, tr(d["foot"], self.lang), top=6.72)

    def s_image(self, d):
        s = self.slide()
        self.heading(s, tr(d["title"], self.lang))
        cap = d.get("caption")
        box = (ML, 1.7, SW - 2 * ML, (4.55 if cap else 5.2))
        self.fit_image(s, os.path.join(ASSETS, d["image"]), box)
        if cap:
            self.caption(s, tr(cap, self.lang), top=6.45)

    def s_two_col(self, d):
        s = self.slide()
        self.heading(s, tr(d["title"], self.lang))
        colw = (SW - 2 * ML - 0.5) / 2
        cols = [d["left"], d["right"]]
        if self.rtl:
            cols = cols[::-1]
        xs = [ML, ML + colw + 0.5]
        for x, col in zip(xs, cols):
            self.rect(s, x, 1.75, colw, 4.95, fill=SURF, line=LINE, lw=1.0,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = self.tb(s, x + 0.32, 2.02, colw - 0.64, 0.5)
            self.para(tf, [(tr(col["head"], self.lang), CYAN, True, 18)], size=18, first=True,
                      before=0, after=0)
            inner = colw - 0.64
            ileft = (x + 0.32)
            tf2 = self.tb(s, ileft, 2.62, inner, 3.9)
            for i, it in enumerate(col["bullets"]):
                txt = tr(it, self.lang)
                runs = [("•  ", CYAN, True, 14.5), (txt, INK, False, 14.5)]
                self.para(tf2, runs, size=14.5, before=(0 if i == 0 else 7), after=0,
                          line=1.12, first=(i == 0))

    def s_table(self, d):
        s = self.slide()
        self.heading(s, tr(d["title"], self.lang))
        head = list(d["head"]); rows = [list(r) for r in d["rows"]]
        if self.rtl:
            head = head[::-1]; rows = [r[::-1] for r in rows]
        ncol = len(head); nrow = len(rows) + 1
        top = 1.9
        tw = SW - 2 * ML
        th = min(4.4, 0.62 * nrow)
        gshape = s.shapes.add_table(nrow, ncol, Inches(ML), Inches(top), Inches(tw), Inches(th))
        tbl = gshape.table
        tbl.first_row = False; tbl.horz_banding = False
        # column widths: give the last (description) column more room
        if ncol == 3:
            widths = [0.12, 0.30, 0.58]
        elif ncol == 2:
            widths = [0.30, 0.70]
        else:
            widths = [1.0 / ncol] * ncol
        if self.rtl:
            widths = widths[::-1]
        for j, frac in enumerate(widths):
            tbl.columns[j].width = Inches(tw * frac)
        for j, htext in enumerate(head):
            self._cell(tbl.cell(0, j), tr(htext, self.lang), bold=True,
                       fg=INKS, bgc=SURF2, size=14.5)
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                self._cell(tbl.cell(i, j), tr(val, self.lang), bold=(j == (ncol - 1 if self.rtl else 0)),
                           fg=INK, bgc=(SURF if i % 2 else BG), size=13)
        if d.get("foot"):
            self.caption(s, tr(d["foot"], self.lang), top=top + th + 0.2)

    def _cell(self, cell, text, *, bold, fg, bgc, size):
        cell.fill.solid(); cell.fill.fore_color.rgb = bgc
        cell.margin_left = cell.margin_right = Inches(0.12)
        cell.margin_top = cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if self.rtl else PP_ALIGN.LEFT
        _set_rtl(p, self.rtl)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = fg
        r.font.name = self.latin
        if self.rtl:
            _set_cs(r, self.arface)

    def s_bigstat(self, d):
        s = self.slide()
        self.rect(s, 0, 0, SW, SH, fill=BG)
        self.accent_bar(s, SW / 2 - 1.0, 1.5, w=2.0, color=CYAN)
        tf = self.tb(s, 0, 1.7, SW, 1.9, anchor=MSO_ANCHOR.MIDDLE)
        self.para(tf, [(d["stat"], CYAN, True, 110)], size=110, first=True, before=0, after=0,
                  align=PP_ALIGN.CENTER)
        tf = self.tb(s, 0, 3.7, SW, 0.6)
        self.para(tf, [(tr(d["label"], self.lang), MUT, False, 20)], size=20, first=True,
                  before=0, after=0, align=PP_ALIGN.CENTER)
        tf = self.tb(s, ML, 4.5, SW - 2 * ML, 0.7)
        self.para(tf, [(tr(d["title"], self.lang), AMB, True, 24)], size=24, first=True,
                  before=0, after=0, align=PP_ALIGN.CENTER)
        tf = self.tb(s, 1.6, 5.3, SW - 3.2, 1.6)
        self.para(tf, [(tr(d["caption"], self.lang), INK, False, 16)], size=16, first=True,
                  before=0, after=0, align=PP_ALIGN.CENTER, line=1.25)

    def s_closing(self, d):
        s = self.slide()
        self.rect(s, -1, 5.7, 15, 3, fill=SURF)
        self.accent_bar(s, (SW - ML - 2.0) if self.rtl else ML, 2.4, w=2.0, color=CYAN)
        tf = self.tb(s, ML, 2.6, SW - 2 * ML, 1.4)
        self.para(tf, [(tr(d["title"], self.lang), INKS, True, 52)], size=52, first=True, before=0, after=0)
        tf = self.tb(s, ML, 4.1, SW - 2 * ML, 1.6)
        self.para(tf, [(tr(d["subtitle"], self.lang), INK, False, 19)], size=19, first=True,
                  before=0, after=0, line=1.25)
        tf = self.tb(s, ML, 6.55, SW - 2 * ML, 0.6)
        self.para(tf, [(tr(d["foot"], self.lang), CYAN, True, 15)], size=15, first=True, before=0, after=0)

    def build(self):
        dispatch = {
            "title": self.s_title, "section": self.s_section, "bullets": self.s_bullets,
            "image": self.s_image, "two_col": self.s_two_col, "table": self.s_table,
            "bigstat": self.s_bigstat, "closing": self.s_closing,
        }
        for d in SLIDES:
            dispatch[d["type"]](d)
        return self.prs


def main():
    for lang, fname in [("en", "NeuroScan_EN.pptx"), ("ar", "NeuroScan_AR.pptx")]:
        deck = Deck(lang)
        prs = deck.build()
        path = os.path.join(OUT, fname)
        prs.save(path)
        print(f"wrote {fname}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
